from .base import BaseParser
from typing import Dict, Any

class PptxParser(BaseParser):
    def _extract(self, file_path: str) -> Dict[str, Any]:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_content = []
        pages = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            joined_slide_text = "\n".join(slide_text)
            text_content.append(joined_slide_text)
            pages.append({"page_num": i + 1, "text": joined_slide_text})
        
        return {
            "text": "\n".join(text_content),
            "pages": pages,
            "tables": [],
            "images": [],
            "metadata": {"source": file_path},
            "sections": []
        }
